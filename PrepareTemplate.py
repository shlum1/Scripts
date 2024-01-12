import collections
import collections.abc
from pptx import Presentation
from sys import argv
from shutil import copy2 as cp
import glob


sVersion = "PrepareTemplate v. 2024-01-12.007"

sYear = '2024'

print('#### ', sVersion)



sDiPPS_DataPath = "//di-pps/dipps_daten/Archiv"

def CopyData(month: str):
    try: 
        if month=='12':
            month = 0
        sMonth=f'{int(month)+1:02}'
    except:
        sMonth = '01'

    print(f' --> CopyData({sYear}-{sMonth}) ...')
    cp(f"//di-pps/dipps_daten/Archiv/DI/DI_{sYear}_{sMonth}_KA_Nettowerte_Pos.csv", "Data-IN")
    cp(f"//di-pps/dipps_daten/Archiv/DI/DI_{sYear}_{sMonth}_Rechnungen.csv", "Data-IN")
    cp(f"//di-pps/dipps_daten/Archiv/TTE/TTE_{sYear}_{sMonth}_KA_Nettowerte_Pos.csv", "Data-IN")
    cp(f"//di-pps/dipps_daten/Archiv/TTE/TTE_{sYear}_{sMonth}_Rechnungen.csv", "Data-IN")

    #for fn in glob.glob(f"//di-daten/verwaltung/GL/TSS/MBR/{sYear}/{month:02}/Data-IN/*.csv"):
    #    cp(fn, "Data-IN")
    for fn in glob.glob(f"//di-daten/verwaltung/GL/TSS/MBR/{sYear}/{month:02}/Data-IN/*.xlsx"):
        cp(fn, "Data-IN")

#******************************************************************************************************************
#*
#*
#*
#******************************************************************************************************************
def FixMonthTitle():
    print('  --> FixMonthTitle() ...')
    
    slide=prs.slides[0]   # title slide

    for s in slide.shapes:   #find text 'Monthly Business Review \n MM / YYYY'
        if s.shape_type==17 and s.text.find('Monthly') >= 0:  #text frame
           break

    tf = s.text_frame
    lp = len(tf.paragraphs)
    p0 = tf.paragraphs[0]
    r = p0.runs
    lr = len(r)
    r1 = r[1]  #2nd paragraph
    r1.text = f'{sMonth} / {sYear}'
    if lr > 2:
        r[2].text = ''



###################################################################################################################
if len(argv) >= 3:
    path = argv[1] 
    sMonth = argv[2]
    try:
        gYear = argv[3]
    except:
        pass  #use pre-defined year
else:
    fName = 'xx/TTE-MBR-2022-xx.pptx'
    sMonth = 'xx'


try:
    fName=f'{path}/TTE-MBR-{sYear}-{sMonth}_leer.pptx'
    print(fName, sYear, sMonth)
    #print(argv)

    prs=Presentation(fName)


    FixMonthTitle()

    #do other stuff to prepare...

    prs.save(fName)

    CopyData(sMonth)


    print('... PrepareTemplate -done- \n\n')
    exit(0)

except Exception as e:
    print("******** ERROR in PrepareTemplate *********")
    print(e)          

    exit(-1)