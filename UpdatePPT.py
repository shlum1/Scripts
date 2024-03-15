import collections
import collections.abc

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
from pptx.util import Inches
import numpy as np 
import datetime


import plotly.graph_objects as go
import plotly.express as px
import pandas as pd

from sys import argv


sVersion = "UpdatePPT v. 2024-03-15.006"
gYear = 2023
#gYear = 2024



#----------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
#
#
#----------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
def DoInvoiceImg(slide, sComp, slidename):
    try:
        print(f'DoInvoiceImg {sComp}')
        slide.shapes.add_picture("Invoices-DI.png", Inches(0.5), Inches(1), height=Inches(5.9), width=Inches(5.9))
        slide.shapes.add_picture("Invoices-TTE.png", Inches(7),  Inches(1), height=Inches(5.9), width=Inches(5.9))
        return True

    except Exception as e:
        print (f" Exception in DoInvoiceImg({slidename}, {sComp}): {e}")     
        return False




#----------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
#
'''
slideNames = {   #2024-03-15.006
    'Title:': None, #1
    'BU Analytics:Month:TTE.': None, #2
    'BU Analytics:Month:DI': None, #3
    'BU Analytics:Month:sum': None, #4
    'BU Analytics:Year:cons.': None, #5
    'BU Analytics:Balance:sum': None, #6
    'Sales:Offers History:DI': None, #7
    'Sales:Offers History:TTE': None, #8
    'Sales:Biggest Offers:DI': None, #9
    'Sales:Biggest Offers:TTE': None, #10

        'Sales:Won and Lost:': None, #11
    'VBP:Ongoing efforts:': None, #12
    'WIP:Orders Backlog:DI': None, #13
    'WIP:Orders Backlog:TTE': None, #14
    'Special Topics:': None, #15
    'INDEN Initiative:': None, #16
    'Finance:Invoices:DI': DoInvoiceImg, #17
    'Finance:Invoices:TTE': DoInvoiceImg, #18
    'R&D Roadmap:DI': None, #19
    'R&D Roadmap:TTE': None,#20
    'HR, Sick leave:': None,#21
    'Thanks:': None #22
} 
'''
slideNames = {   #2024-03-15.006
    'Title:': None, #1
    'BU Analytics:Month:TTE.': None, #2
    'BU Analytics:Month:DI': None, #3
    'BU Analytics:Month:sum': None, #4
    'BU Analytics:Year:cons.': None, #5
    'BU Analytics:Balance:sum': None, #6
    'Sales:Offers History:DI': None, #7
    'Sales:Offers History:TTE': None, #8
    'Sales:Biggest Offers:DI': None, #9
    'Sales:Biggest Offers:TTE': None, #10
    'Sales:Special Topics:DI': None, #11
    'Sales:Special Topics:TTE': None, #12
    'VBP:Status VBP:DI:': None, #13
    'VBP:Status VBP:DI:': None, #14
    'Sales:Attrition:': None, #15
    'PS:WIP:DI': None, #16
    'PS:WIP:TTE': None, #17
    'CS:Tickets:TTE': None,#18
    'HR:Sick leave:': None,#19
    'Thanks:': None #20
} 





###################################################################################################################################################################################
###################################################################################################################################################################################
if __name__ == '__main__':
    print('#### ', sVersion)

    if len(argv) >= 3:
        path = argv[1]
        sMonth = argv[2]
        try:
            gYear = argv[3]
        except:
            pass  #use pre-defined year
    else:
        sMonth = 'xx'
        path = f'test/{sMonth}'
        print("Usage: python UpdatePPT.py <path> <month>")
        #exit(1)

    
    try:
        fName = f'{path}/TTE-MBR-{gYear}-{sMonth}_leer.pptx' 
        #print(argv)

        print(fName, gYear, sMonth)

        prs = Presentation(fName)

        # update all slides if necessary...
        for i, sSlide in enumerate(slideNames):
            fnc = slideNames[sSlide]
            slide = prs.slides[i]    
            print (f'Slide {i} : {sSlide}')
            l = sSlide.split(':')
            if fnc != None:
                fnc(slide, l[-1:][0], sSlide)

        

        prs.save(f'{path}/TTE-MBR-{gYear}-{sMonth}-Draft.pptx') # saving file
        
        print('... UpdatePPT -done- \n\n')

        exit(0)

    except Exception as e:
        print("******** ERROR in UpdatePPT *********")
        print(e)          

        exit(-1)      







